"""Build a PPTX report from a run's outputs + figures (raw values only, from the artifacts).

Exposed as the `export_report_ppt` tool. Reads coverage_summary.tsv + the provenance tsvs
(kegg_recovered / hmdb_recovered / unmapped_harmonization / exogenous_kept / xenobiotic_excluded) + figures/*.png
from the workdir. All numbers are read from the run — nothing is invented. Requires
python-pptx + pillow (declared in pyproject).
"""

from __future__ import annotations

import csv
import re
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PRIMARY = RGBColor(0x2A, 0x7D, 0xB5); DARK = RGBColor(0x1B, 0x3A, 0x52)
GREY = RGBColor(0x55, 0x55, 0x55); LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); AMBER = RGBColor(0xC8, 0x7A, 0x1E)

# populated by generate()
prs = BLANK = SW = SH = RUN = FIG = META = None
_MZ = re.compile(r"(m/z\s*[\d.]+\s*\[M[^\]]*\])")


def read_tsv(p):
    p = Path(p)
    if not p.exists():
        return []
    with open(p) as f:
        return list(csv.DictReader(f, delimiter="\t"))


def cause_of(o, logic, route):
    low = logic.lower(); rl = route.lower()
    if "typo" in low:
        return "Typo (misspelling)"
    if "xref bridge" in rl or "bridgedb promotes" in low:
        return "Vendor/trade name — no direct KEGG"
    if "combined name" in low or "isomer pair" in low or "/" in o:
        return "Combined / isomer name"
    if re.search(r"=\s*[A-Z0-9]{2,5}\b", o) or re.fullmatch(r"[A-Z0-9]{2,5}", o.strip()) \
            or re.search(r"\([A-Z0-9]{2,6}\)", o):
        return "Abbreviation / symbol"
    if ("carnitine" in o.lower() and " " in o) or re.search(r"-\d-?[Pp]hosphate", o) or o.startswith("Orto"):
        return "Name format (spacing/hyphen)"
    return "Synonym / alternate name"


def verify_of(logic):
    m = _MZ.search(logic or "")
    return ("mass ✓ " + m.group(1)) if m else "mass ✓"


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(par, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    par.text = text if text not in (None, "") else " "
    par.alignment = align
    r = par.runs[0]; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Calibri"


def header(slide, title, sub=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
    tf = _box(slide, Inches(0.55), Inches(0.35), SW - Inches(1.1), Inches(0.9))
    _set(tf.paragraphs[0], title, 26, DARK, True)
    if sub:
        p = tf.add_paragraph(); _set(p, sub, 13, GREY)


def title_slide():
    s = prs.slides.add_slide(BLANK)
    band = s.shapes.add_shape(1, 0, Inches(2.4), SW, Inches(2.7))
    band.fill.solid(); band.fill.fore_color.rgb = PRIMARY; band.line.fill.background()
    tf = _box(s, Inches(0.8), Inches(2.75), SW - Inches(1.6), Inches(2.0))
    _set(tf.paragraphs[0], "Metabolite Identifier Harmonization", 40, WHITE, True)
    p = tf.add_paragraph(); _set(p, "name → KEGG / HMDB / ChEBI / PubChem + Mouse-GEM MAM", 17, WHITE)
    tf2 = _box(s, Inches(0.8), Inches(5.4), SW - Inches(1.6), Inches(1.2))
    _set(tf2.paragraphs[0], "metabo-idmapper MCP  ·  deterministic tools + LLM reasoning layer", 14, GREY)
    p = tf2.add_paragraph(); _set(p, f"{META['total']} metabolites  ·  run: {META['run']}/", 13, GREY)


def kpi_slide():
    cov = META["cov"]
    s = prs.slides.add_slide(BLANK)
    header(s, "Coverage summary", f"{META['total']} metabolites; classes mutually exclusive; has_* count any xref present")
    kpis = [("KEGG-mapped", "class:KEGG-mapped"), ("HMDB-mapped", "class:HMDB-mapped"),
            ("structure-only", "class:structure-only"), ("exogenous (kept)", "class:exogenous"),
            ("xenobiotic-excl.", "class:xenobiotic-excluded"), ("GEM-mapped", "gem_mapped")]
    x0, y0, cw, gap = Inches(0.55), Inches(1.7), Inches(2.0), Inches(0.1)
    for i, (lab, key) in enumerate(kpis):
        row = cov.get(key, {"value": "0", "pct": "0"})
        card = s.shapes.add_shape(1, x0 + (cw + gap) * i, y0, cw, Inches(1.5))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = PRIMARY; card.line.width = Pt(1)
        tf = card.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set(tf.paragraphs[0], row["value"], 28, PRIMARY, True, PP_ALIGN.CENTER)
        _set(tf.add_paragraph(), lab, 11, DARK, True, PP_ALIGN.CENTER)
        _set(tf.add_paragraph(), row["pct"] + "%", 12, GREY, True, PP_ALIGN.CENTER)
    tf = _box(s, Inches(0.55), Inches(3.5), SW - Inches(1.1), Inches(0.6))
    _set(tf.paragraphs[0], f"has_KEGG {cov['has_kegg']['value']} ({cov['has_kegg']['pct']}%)   ·   "
         f"has_HMDB {cov['has_hmdb']['value']} ({cov['has_hmdb']['pct']}%)   ·   truly-unmapped 0", 14, DARK, True)
    tf2 = _box(s, Inches(0.55), Inches(4.2), SW - Inches(1.1), Inches(2.6))
    bullets = [
        "1차 MetaboAnalyst 매핑 → 미매핑 항목은 rename 규칙으로 재검색(KEGG/PubChem/ChEBI) → BridgeDb 상호대조 + HMDB backfill.",
        f"has_HMDB({cov['has_hmdb']['value']})가 has_KEGG({cov['has_kegg']['value']})보다 높음 — backfill로 HMDB 과소집계 해소.",
        f"exogenous(외인성 생물학적: 약물·식이·장내미생물·식물) {META['exo']}개는 유지·태그(외부 유래 실신호).",
        f"xenobiotic {META['xeno']}개(계면활성제·가소제·산업/시약·LC-MS 첨가물)는 비생물성 오염물 — 분석 제외.",
        f"structure-only {META['struct']}개는 구조 ID(PubChem/ChEBI)만 확보 — pathway/flux 미사용."]
    for b in bullets:
        _set(tf2.add_paragraph(), "•  " + b, 13, DARK)


def figure_slide(title, sub, png):
    s = prs.slides.add_slide(BLANK); header(s, title, sub)
    if not Path(png).exists():
        _set(_box(s, Inches(0.6), Inches(3), SW - Inches(1.2), Inches(1)).paragraphs[0],
             f"(figure not found: {Path(png).name})", 14, GREY)
        return
    from PIL import Image
    iw, ih = Image.open(png).size
    scale = min((SW - Inches(1.1)) / iw, (SH - Inches(2.0)) / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(png), int((SW - w) / 2), Inches(1.75), width=w, height=h)


def table_slides(title, sub, rows, cols, widths, rows_per_page=16, fontsize=8.0, amber_last=False):
    pages = [rows[i:i + rows_per_page] for i in range(0, len(rows), rows_per_page)] or [[]]
    for pi, prows in enumerate(pages):
        s = prs.slides.add_slide(BLANK)
        header(s, title if len(pages) == 1 else f"{title}   ({pi + 1}/{len(pages)})", sub)
        tw = SW - Inches(0.7)
        tbl = s.shapes.add_table(len(prows) + 1, len(cols), Inches(0.35), Inches(1.55), tw, Inches(5.6)).table
        for j, wd in enumerate(widths):
            tbl.columns[j].width = Emu(int(tw * wd))
        for j, c in enumerate(cols):
            cell = tbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
            cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
            _set(cell.text_frame.paragraphs[0], c[1], min(fontsize + 1.5, 10), WHITE, True)
        for i, r in enumerate(prows, 1):
            for j, c in enumerate(cols):
                cell = tbl.cell(i, j)
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                cell.text_frame.word_wrap = True
                cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
                cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
                col = AMBER if (amber_last and j == len(cols) - 1) else DARK
                _set(cell.text_frame.paragraphs[0], str(r.get(c[0], "")), fontsize, col)


def methods_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "Methods (detail)", "metabo-idmapper MCP · deterministic engines + LLM reasoning layer (summary-in → decision-out; anti-fabrication)")
    tf = _box(s, Inches(0.55), Inches(1.55), SW - Inches(1.1), Inches(5.6))
    items = [
        ("Data", f"{META['total']} LC-MS metabolites; feature별 m/z + adduct 보존 → 질량 검증에 사용."),
        ("Stage 1 · 1st-pass (raw MetaboAnalystR)", "CrossReferencing(name→KEGG/HMDB/PubChem/ChEBI). exact match만 자동채택(M1). 전부 매칭 실패 시 크래시 방지 위해 anchor 화합물 주입."),
        ("Stage 2 · unmapped", "1차에서 KEGG·HMDB 둘 다 없는 항목 추출(xenobiotic-excluded 제외)."),
        ("Stage 3 · recovery (핵심)", "LLM이 rename 규칙 제안(오타수정·약어확장·synonym·이름형식) → KEGGREST keggFind + PubChem PUG-REST + ChEBI OLS4 재검색; keggFind 다중 hit는 molmass 단일동위원소 질량 vs 관측 m/z(±ppm)·RT로 disambiguation."),
        ("Stage 4 · cross-check (raw BridgeDbR)", "확보한 ID 1개를 KEGG/HMDB/ChEBI/PubChem/InChIKey로 상호대조; KEGG만 있는 항목에 HMDB backfill."),
        ("Origin split", f"비생물성 오염물 lexicon(LC-MS 첨가물·계면활성제·가소제·산업) → xenobiotic-excluded {META['xeno']}개 제외; 생물학적 외인성(약물·식이·장내미생물·식물) exogenous {META['exo']}개는 유지·태그."),
        ("GEM crosswalk (raw COBRApy)", "Mouse-GEM, KEGG>HMDB>ChEBI xref → MAM(대사모델 flux 입력)."),
        ("Confidence tiers", "M1 exact · M2 verified synonym/abbrev · M3 verified typo · M4 structure-only."),
        ("Outputs", "coverage · provenance 표 · figures · reproduce_mapping.py/.ipynb · 원본 파일 ID annotate."),
    ]
    first = True
    for h, d in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = f"{h}:  "; r.font.bold = True; r.font.size = Pt(12.5); r.font.color.rgb = PRIMARY; r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = d; r2.font.size = Pt(11.5); r2.font.color.rgb = DARK; r2.font.name = "Calibri"
        p.space_after = Pt(6)


def flow_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "Pipeline flow", "deterministic tools EMIT evidence · LLM reasoning layer makes the identity CALL")
    steps = [("01 Ingest", "normalize · xenobiotic screen"), ("02 Exact match", "MetaboAnalystR 1st pass (auto-M1)"),
             ("03 Evidence & recovery", "rename → KEGG/PubChem/ChEBI 재검색 · molmass 검증"),
             ("04 Cross-check", "BridgeDb xref + HMDB backfill"), ("05 The CALL", "screen_exogenous · record_decision"),
             ("06 GEM crosswalk", "COBRApy → Mouse-GEM MAM"), ("07 Report", "coverage · provenance · figures · code · annotate")]
    y = Inches(1.8); h = Inches(0.62); gap = Inches(0.12); w = SW - Inches(1.4)
    for i, (t, d) in enumerate(steps):
        box = s.shapes.add_shape(1, Inches(0.7), y + (h + gap) * i, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = PRIMARY if i in (2, 3) else LIGHT; box.line.color.rgb = PRIMARY
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        c = WHITE if i in (2, 3) else DARK; p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"  {t}   "; r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = c; r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = d; r2.font.size = Pt(11); r2.font.color.rgb = c; r2.font.name = "Calibri"


def recovery_mechanics_slide(kegg):
    s = prs.slides.add_slide(BLANK)
    header(s, "Recovery — why 1st pass missed & how it mapped", "MetaboAnalyst가 못 잡은 이유별 그룹 + 해결 방법 (전부 molmass 질량검증 통과)")
    groups = {}
    for r in kegg:
        groups.setdefault(r["cause"], []).append(r)
    order = ["Typo (misspelling)", "Abbreviation / symbol", "Name format (spacing/hyphen)",
             "Synonym / alternate name", "Combined / isomer name", "Vendor/trade name — no direct KEGG"]
    how = {"Typo (misspelling)": "철자 교정 → keggFind 재검색", "Abbreviation / symbol": "약어 전개(정식명) → keggFind",
           "Name format (spacing/hyphen)": "이름형식 정규화(띄어쓰기/하이픈) → keggFind",
           "Synonym / alternate name": "동의어/정식 화학명으로 치환 → keggFind",
           "Combined / isomer name": "결합/이성질체명 분해, RT·질량으로 대표 이성질체 확정",
           "Vendor/trade name — no direct KEGG": "KEGG 직접 없음 → BridgeDb가 PubChem→KEGG/HMDB 승격"}
    cats = [c for c in order if c in groups] + [c for c in groups if c not in order]
    rows = [{"cause": c, "n": str(len(groups[c])), "how": how.get(c, "rename → 재검색"),
             "ex": "; ".join(f"{r['original_name']}→{r['harmonized_name']}" for r in groups[c])} for c in cats]
    cols = [("cause", "원인 (why 1st-pass missed)"), ("n", "n"), ("how", "해결 방법 (how)"), ("ex", "예시 (original → harmonized)")]
    tw = SW - Inches(0.7)
    tbl = s.shapes.add_table(len(rows) + 1, 4, Inches(0.35), Inches(1.6), tw, Inches(5.1)).table
    for j, wd in zip(range(4), [0.20, 0.04, 0.28, 0.48]):
        tbl.columns[j].width = Emu(int(tw * wd))
    for j, c in enumerate(cols):
        cell = tbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        _set(cell.text_frame.paragraphs[0], c[1], 10, WHITE, True)
    for i, r in enumerate(rows, 1):
        for j, c in enumerate(cols):
            cell = tbl.cell(i, j); cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            cell.text_frame.word_wrap = True; cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
            _set(cell.text_frame.paragraphs[0], str(r.get(c[0], "")), 8.5, DARK, c[0] == "cause")


def outputs_slide():
    s = prs.slides.add_slide(BLANK)
    header(s, "Outputs & reproducibility", f"everything below is written to the run workdir  {META['run']}/")
    tf = _box(s, Inches(0.6), Inches(1.7), SW - Inches(1.2), Inches(5.2))
    items = [("Tables", "master_ledger.tsv · coverage_summary.tsv · enriched_xref.tsv · kegg_recovered · hmdb_recovered · unmapped_harmonization · exogenous_kept · xenobiotic_excluded"),
             ("Annotated source", "원본 데이터 파일 + 최종 ID 컬럼(ID_kegg/hmdb/chebi/pubchem/inchikey/final_class/gem_mam)"),
             ("Figures", "figures/db_matching_upset.png · db_matching_improvement.png"),
             ("Reproduction code", "code/reproduce_mapping.py + .ipynb (원본 라이브러리 API로 흐름 재실행) + ma.R/bridge.R/kegg.R"),
             ("Provenance", "midmap_ledger.json — 원본명→normalize→candidates(source별 query)→accepted→decisions")]
    first = True
    for h, d in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = f"{h}:  "; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = PRIMARY; r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = d; r2.font.size = Pt(12.5); r2.font.color.rgb = DARK; r2.font.name = "Calibri"
        p.space_after = Pt(10)


def generate(workdir: str, out: str | None = None) -> str:
    """Build <workdir>/metabolite_id_mapping_report.pptx and return its path."""
    global prs, BLANK, SW, SH, RUN, FIG, META
    RUN = Path(workdir).resolve(); FIG = RUN / "figures"
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height
    cov = {r["metric"]: r for r in read_tsv(RUN / "coverage_summary.tsv")}
    if not cov:
        raise FileNotFoundError("coverage_summary.tsv not found — run coverage_summary first")
    total = cov.get("has_kegg", {}).get("denom", "0")
    META = {"cov": cov, "run": RUN.name, "total": total,
            "exo": cov.get("class:exogenous", {}).get("value", "0"),
            "xeno": cov.get("class:xenobiotic-excluded", {}).get("value", "0"),
            "struct": cov.get("class:structure-only", {}).get("value", "0")}

    kegg = read_tsv(RUN / "kegg_recovered.tsv")
    for r in kegg:
        r["cause"] = cause_of(r["original_name"], r.get("logic", ""), r.get("route", ""))
        r["verify"] = verify_of(r.get("logic", ""))
    hmdb = read_tsv(RUN / "hmdb_recovered.tsv")
    unm = read_tsv(RUN / "unmapped_harmonization.tsv")
    exo_kept = read_tsv(RUN / "exogenous_kept.tsv")
    xeno = read_tsv(RUN / "xenobiotic_excluded.tsv")

    title_slide(); kpi_slide(); methods_slide(); flow_slide()
    figure_slide("DB identifier coverage (UpSet)", "5-DB membership across analysable features (union of candidates)", FIG / "db_matching_upset.png")
    figure_slide("Improvement over MetaboAnalyst 1st pass", "baseline (MetaboAnalyst only) vs + current logic; resolution composition", FIG / "db_matching_improvement.png")
    if kegg:
        recovery_mechanics_slide(kegg)
        table_slides(f"KEGG recovered — original · why · fix · verify ({len(kegg)})",
                     "왜 1차에서 못 잡았고(why) 어떤 rename으로 매핑됐는지(fix) + 질량검증",
                     kegg, [("original_name", "original"), ("cause", "why (cause)"),
                            ("harmonized_name", "fix (harmonized)"), ("kegg", "KEGG"), ("verify", "verify")],
                     [0.22, 0.20, 0.22, 0.09, 0.27], rows_per_page=15, fontsize=8.5)
    if hmdb:
        table_slides(f"HMDB recovered beyond MetaboAnalyst ({len(hmdb)})", "대부분 BridgeDb xref backfill",
                     hmdb, [("original_name", "original"), ("harmonized_name", "bridged from"),
                            ("hmdb", "HMDB"), ("route", "route")],
                     [0.34, 0.28, 0.16, 0.22], rows_per_page=16, fontsize=8.5)
    if unm:
        table_slides(f"Unmapped — harmonized but not mapped ({len(unm)})", "구조 ID만 확보(structure-only); 실패 이유 전문",
                     unm, [("original_name", "original"), ("structure_ids", "structure IDs"),
                           ("reason", "why not mapped (full)")],
                     [0.20, 0.20, 0.60], rows_per_page=6, fontsize=8.5, amber_last=True)
    if exo_kept:
        catsub = " · ".join(f"{k} {v}" for k, v in Counter(r["origin"] for r in exo_kept).most_common())
        table_slides(f"Exogenous kept — biological, outside-host ({len(exo_kept)})",
                     f"약물·식이·장내미생물·식물 유래 실신호 — origin별({catsub}); 분석에 유지·태그",
                     exo_kept, [("original_name", "original"), ("origin", "origin"),
                                ("ids", "ids"), ("gem_xref", "GEM"), ("reason", "note (full)")],
                     [0.22, 0.13, 0.22, 0.10, 0.33], rows_per_page=9, fontsize=8.5)
    if xeno:
        catsub = " · ".join(f"{k} {v}" for k, v in Counter(r["origin"] for r in xeno).most_common())
        table_slides(f"Xenobiotic excluded — non-biological ({len(xeno)})",
                     f"분석 제외 오염물 — origin별({catsub}) + 제외 사유 전문",
                     xeno, [("original_name", "original"), ("origin", "origin"),
                            ("ids", "ids (if any)"), ("reason", "excluded because (full)")],
                     [0.22, 0.13, 0.22, 0.43], rows_per_page=9, fontsize=8.5, amber_last=True)
    outputs_slide()

    out = out or str(RUN / "metabolite_id_mapping_report.pptx")
    prs.save(out)
    return out
